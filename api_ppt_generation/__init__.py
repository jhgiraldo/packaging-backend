import azure.functions as func
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import os
import requests
import matplotlib.pyplot as plt

# URL pública del Blob Storage
TEMPLATE_URL = "https://hibecistorage1.blob.core.windows.net/blob-publico/PPT Generator Template HAVAS.pptx"

def generar_grafica(planned, delivered, categorias, output_path):
    fig, ax = plt.subplots(figsize=(10, 6))
    bar_width = 0.35
    x = range(len(categorias))

    # Barras
    ax.bar(x, planned, width=bar_width, label="Planned Spend", color="#A52A2A")
    ax.bar([i + bar_width for i in x], delivered, width=bar_width, label="Delivered Spend", color="#FFD700")

    # Etiquetas y título
    ax.set_xticks([i + bar_width/2 for i in x])
    ax.set_xticklabels(categorias, rotation=45, ha="right")
    ax.set_ylabel("£")
    ax.set_title("Planned vs Delivered Spend")
    ax.legend()

    # Guardar con fondo transparente
    plt.tight_layout()
    plt.savefig(output_path, transparent=True)
    plt.close(fig)

def main(req: func.HttpRequest) -> func.HttpResponse:
    try:
        # Leer datos del body
        body = req.get_json()
        title = body.get("title", "Informe")
        subtitle = body.get("subtitle", "")
        bullets = body.get("bullets", [
            "Total planned: £4.3M | Delivered: £324K (8%)",
            "Digital pacing steady: OLV and Social platforms delivering 63–78% of planned spend.",
            "PPC ramp-up: Current delivery is between 22–25% as activation is phased.",
            "Gap expected: Digital providing strong early base; PPC will scale further."
        ])
        planned = body.get("planned", [30000,25000,20000,40000,60000,35000,30000,150000,50000,10000])
        delivered = body.get("delivered", [15000,12000,10000,30000,40000,25000,20000,33000,12000,5000])
        categorias = body.get("categorias", ["YouTube View","YouTube Reach","Havas Marketplace","Teads","META","LinkedIn","TikTok","Google Search","Google Demand","Bing"])

        # Descargar plantilla desde Blob Storage
        response = requests.get(TEMPLATE_URL)
        if response.status_code != 200:
            return func.HttpResponse("Error al descargar la plantilla desde Blob Storage", status_code=500)

        prs = Presentation(io.BytesIO(response.content))

        # Generar gráfica transparente
        grafico_path = os.path.join(os.getcwd(), "grafico.png")
        generar_grafica(planned, delivered, categorias, grafico_path)

        # Buscar la marca plannedvsdeliveredtext
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "plannedvsdeliveredtext" in shape.text:
                    # Insertar imagen encima del bloque de texto
                    left = shape.left
                    top = shape.top - Inches(4)  # mover arriba
                    width = Inches(6)
                    height = Inches(4)
                    slide.shapes.add_picture(grafico_path, left, top, width, height)

                    # Reemplazar texto con título y bullets
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p_title = text_frame.add_paragraph()
                    p_title.text = f"{title} - {subtitle}"
                    p_title.font.size = Pt(18)
                    p_title.font.bold = True

                    for bullet in bullets:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(14)
                        p.level = 0
                    break

        # Guardar en memoria
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)

        return func.HttpResponse(
            ppt_stream.read(),
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=planned_vs_delivered.pptx"}
        )

    except Exception as e:
        return func.HttpResponse(f"Error interno: {str(e)}", status_code=500)